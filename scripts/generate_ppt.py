"""
生成金融风格PPT：商业银行代客外汇交易、理财代销业务——客户准入全流程及系统实现方案
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==================== 金融配色方案 ====================
NAVY_DARK = RGBColor(0x0B, 0x1D, 0x3A)      # 深海军蓝（背景）
NAVY = RGBColor(0x14, 0x2D, 0x5E)            # 海军蓝
GOLD = RGBColor(0xC5, 0x9B, 0x3C)            # 金色（强调）
GOLD_LIGHT = RGBColor(0xE8, 0xC8, 0x6A)      # 浅金
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF3, 0xF5)
GRAY = RGBColor(0x8A, 0x8D, 0x93)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
MEDIUM_BLUE = RGBColor(0x1E, 0x4D, 0x8C)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN_ACCENT = RGBColor(0x27, 0xAE, 0x60)
ORANGE_ACCENT = RGBColor(0xE6, 0x7E, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ==================== 辅助函数 ====================
def add_bg(slide, color=NAVY_DARK):
    """添加纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, opacity=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if opacity is not None:
        shape.fill.fore_color.brightness = 0
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=DARK_TEXT,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_name='Microsoft YaHei'):
    """添加多行文本框，lines = [(text, font_size, color, bold, alignment), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text, font_size, color, bold = line[:4]
        align = line[4] if len(line) > 4 else PP_ALIGN.LEFT
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(4)
    return txBox

def add_gold_line(slide, left, top, width):
    """添加金色分隔线"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, content_lines, title_color=NAVY):
    """添加卡片式内容块"""
    # 卡片背景
    card = add_rect(slide, left, top, width, height, WHITE)
    card.shadow.inherit = False
    # 顶部色条
    add_rect(slide, left, top, width, Pt(4), title_color)
    # 标题
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
                 title, font_size=14, color=title_color, bold=True)
    # 内容
    add_multi_text(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4),
                   height - Inches(0.6), content_lines)

def add_page_number(slide, num, total=16):
    """添加页码"""
    add_text_box(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.35),
                 f"{num}/{total}", font_size=9, color=GRAY, alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, section_num, title, subtitle=""):
    """添加章节标题栏"""
    # 顶部色带
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), NAVY)
    add_gold_line(slide, Inches(0.5), Inches(1.1), Inches(12.33))
    # 章节编号
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(1.2), Inches(0.6),
                 f"PART {section_num}", font_size=12, color=GOLD, bold=True)
    # 标题
    add_text_box(slide, Inches(0.5), Inches(0.45), Inches(11), Inches(0.55),
                 title, font_size=24, color=WHITE, bold=True)

# ==================== 第1页：封面 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY_DARK)

# 装饰：左侧金色竖条
add_rect(slide, Inches(0.8), Inches(1.5), Pt(6), Inches(4.5), GOLD)

# 主标题
add_text_box(slide, Inches(1.3), Inches(2.0), Inches(10), Inches(1.0),
             "商业银行代客外汇交易、理财代销业务", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1.3), Inches(2.9), Inches(10), Inches(0.8),
             "客户准入全流程及系统实现方案", font_size=30, color=GOLD, bold=True)

# 分隔线
add_gold_line(slide, Inches(1.3), Inches(3.7), Inches(3))

# 副标题
add_text_box(slide, Inches(1.3), Inches(4.0), Inches(10), Inches(0.5),
             "Customer Onboarding Process & System Implementation", font_size=14, color=GOLD_LIGHT)
add_text_box(slide, Inches(1.3), Inches(4.5), Inches(10), Inches(0.5),
             "涵盖：监管合规 · 准入规则 · 端到端流程 · 系统架构 · 风险管控", font_size=13, color=GRAY)

# 底部
add_text_box(slide, Inches(1.3), Inches(6.3), Inches(5), Inches(0.4),
             "金融市场部  |  合规部  |  风险管理部  |  信息科技部", font_size=10, color=GRAY)
add_text_box(slide, Inches(1.3), Inches(6.6), Inches(5), Inches(0.4),
             "2025 年 5 月  |  V1.0  |  内部限制", font_size=9, color=GRAY)

# 右上角装饰
add_rect(slide, Inches(11.5), Inches(0), Inches(1.83), Inches(7.5), NAVY)

# ==================== 第2页：目录 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 0, "目  录", "CONTENTS")
add_page_number(slide, 2)

toc_items = [
    ("01", "监管合规与准入规则", "适用监管文件 · 个人客户准入条件 · 对公客户准入条件 · 禁止准入清单 · 异常识别规则"),
    ("02", "实际业务操作全流程", "端到端 8 大环节 · 材料清单 · 审批层级 · 异常场景处理"),
    ("03", "系统实现方案", "系统架构 · 功能模块 · 规则引擎 · 数据流转 · 接口设计"),
    ("04", "风险点与管控措施", "合规/操作/洗钱风险 · 三层风控体系 · 监控预警 · 审计留痕"),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.6) + Inches(1.4) * i
    # 编号
    add_text_box(slide, Inches(0.8), y, Inches(0.8), Inches(0.6),
                 num, font_size=32, color=GOLD, bold=True)
    # 竖线
    add_rect(slide, Inches(1.6), y + Inches(0.05), Pt(3), Inches(0.7), GOLD)
    # 标题
    add_text_box(slide, Inches(1.9), y + Inches(0.05), Inches(5), Inches(0.45),
                 title, font_size=20, color=NAVY, bold=True)
    # 描述
    add_text_box(slide, Inches(1.9), y + Inches(0.5), Inches(8), Inches(0.35),
                 desc, font_size=11, color=GRAY)

# ==================== 第3页：监管文件总览 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 1, "监管合规与准入规则 —— 适用监管文件")
add_page_number(slide, 3)

regs = [
    ("《银行业金融机构衍生产品交易\n业务管理暂行办法》", "银监会令 2011年第1号\n（2019年修订）", "客户须具备真实需求背景\n建立客户准入制度\n充分评估风险承受能力"),
    ("《商业银行理财业务监督\n管理办法》", "银保监会令\n2018年第6号", "销售前风险承受能力评估\n不得误导销售\n专区销售+双录要求"),
    ("《关于规范金融机构资产管理\n业务的指导意见》(资管新规)", "银发〔2018〕\n106号", "打破刚兑、净值化管理\n投资者适当性管理\n合格投资者认定"),
    ("《金融机构客户尽职调查和\n客户身份资料保存管理办法》", "央行令\n2022年第1号", "客户尽职调查制度\n识别受益所有人\n持续监控客户交易"),
    ("《中华人民共和国\n反洗钱法》", "2024年修订", "大额和可疑交易报告\n客户身份识别\n制裁名单筛查"),
    ("《国家外汇管理局关于规范\n银行间外汇市场交易的通知》", "外管局", "实需原则\n真实贸易背景材料\n代客结售汇合规要求"),
]

for i, (title, source, req) in enumerate(regs):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(4.1) * col
    y = Inches(1.5) + Inches(2.9) * row
    add_card(slide, x, y, Inches(3.8), Inches(2.6), title, [
        (f"法规来源：{source}", 9, GRAY, False),
        ("", 6, GRAY, False),
        ("核心要求：", 10, NAVY, True),
        (req, 9, DARK_TEXT, False),
    ], title_color=MEDIUM_BLUE)

# ==================== 第4页：个人客户准入条件 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 1, "个人客户准入条件 —— 五维准入模型")
add_page_number(slide, 4)

dims = [
    ("身份条件", NAVY, [
        "年满18周岁，完全民事行为能力",
        "有效二代身份证（联网核查+人脸识别）",
        "已完成I类/II类账户开立",
        "实名认证通过",
    ]),
    ("风险测评", MEDIUM_BLUE, [
        "首次办理前强制完成《风险承受能力评估问卷》",
        "测评内容：年龄、收入、投资经验、风险偏好、亏损承受意愿",
        "有效期：1年，到期前1个月提醒更新",
        "5级：C1保守型 → C5进取型",
    ]),
    ("适当性匹配", GREEN_ACCENT, [
        "客户等级 ≥ 产品等级方可准入",
        "C1级禁止销售III类及以上产品",
        "衍生品：金融资产≥50万或相关交易经验",
        "冷静期：理财认购后24h可撤销",
    ]),
    ("反洗钱筛查", ORANGE_ACCENT, [
        "公安部在逃人员名单",
        "联合国安理会/OFAC SDN制裁名单",
        "央行反洗钱监测名单 + 行内黑名单",
        "PEPS筛查：政治公众人物须EDD",
    ]),
    ("交易权限", RED_ACCENT, [
        "系统根据风险等级自动赋权",
        "超风险品种系统自动拦截",
        "客户可主动关闭品种权限",
        "权限变更需走审批流程",
    ]),
]

for i, (title, color, items) in enumerate(dims):
    x = Inches(0.35) + Inches(2.55) * i
    y = Inches(1.5)
    # 顶部色条
    add_rect(slide, x, y, Inches(2.35), Pt(4), color)
    # 标题
    add_text_box(slide, x, y + Inches(0.1), Inches(2.35), Inches(0.35),
                 title, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # 条目
    lines = [(f"• {item}", 10, DARK_TEXT, False) for item in items]
    add_multi_text(slide, x + Inches(0.1), y + Inches(0.55), Inches(2.15), Inches(2.5), lines)

# 底部风险等级条
grades = [("C1 保守型", LIGHT_GRAY), ("C2 稳健型", RGBColor(0xBE, 0xD9, 0xEA)), ("C3 平衡型", RGBColor(0x8E, 0xBB, 0xD3)),
          ("C4 成长型", MEDIUM_BLUE), ("C5 进取型", NAVY)]
add_text_box(slide, Inches(0.35), Inches(5.1), Inches(4), Inches(0.3),
             "客户风险等级（五级）→", font_size=10, color=GRAY, bold=True)
for i, (label, color) in enumerate(grades):
    x = Inches(0.35) + Inches(2.55) * i
    add_rect(slide, x, Inches(5.4), Inches(2.35), Inches(0.35), color)
    add_text_box(slide, x, Inches(5.4), Inches(2.35), Inches(0.35),
                 label, font_size=10, color=WHITE if i >= 3 else DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

# 产品分界线
prod_grades = [("I类\n低风险", LIGHT_GRAY), ("II类\n中低风险", RGBColor(0xBE, 0xD9, 0xEA)), ("III类\n中风险", RGBColor(0x8E, 0xBB, 0xD3)),
               ("IV类\n中高风险", MEDIUM_BLUE), ("V类\n高风险", NAVY)]
add_text_box(slide, Inches(0.35), Inches(6.0), Inches(4), Inches(0.3),
             "产品风险等级（五级）→", font_size=10, color=GRAY, bold=True)
for i, (label, color) in enumerate(prod_grades):
    x = Inches(0.35) + Inches(2.55) * i
    add_rect(slide, x, Inches(6.3), Inches(2.35), Inches(0.55), color)
    add_multi_text(slide, x, Inches(6.3), Inches(2.35), Inches(0.55),
                   [(label, 9, WHITE if i >= 3 else DARK_TEXT, False, PP_ALIGN.CENTER)])

# ==================== 第5页：对公客户准入条件 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 1, "对公客户准入条件 —— 五维准入 + 行业禁入")
add_page_number(slide, 5)

# 左侧五维
ent_dims = [
    ("资质要求", "有效存续企业法人\n持有营业执照+公司章程\n经营范围含外汇/理财业务\n已完成企业银行账户开立"),
    ("授权材料", "法定代表人身份证明\n授权委托书(非法人办理)\n被授权人身份证明\n交易权限申请书"),
    ("风险等级评定", "6维评估:基本信息/财务/信用\n交易经验/风控制度/实控人\n4级:E1成熟→E4严格限制\n近3年审计报告为必须项"),
    ("反洗钱尽调", "受益所有人穿透至自然人\n(持股≥25%或实际控制人)\n企业股权结构图\n对公制裁筛查+贸易背景核查"),
    ("行业禁入", "房地产(限制类)\n两高一剩(钢铁/煤炭/水泥等)\n政府融资平台\n虚拟货币/博彩行业\n严重违法失信企业"),
]

for i, (title, content) in enumerate(ent_dims):
    y = Inches(1.5) + Inches(1.1) * i
    colors = [NAVY, MEDIUM_BLUE, GREEN_ACCENT, ORANGE_ACCENT, RED_ACCENT]
    add_rect(slide, Inches(0.5), y, Pt(4), Inches(0.9), colors[i])
    add_text_box(slide, Inches(0.7), y, Inches(1.8), Inches(0.3),
                 title, font_size=12, color=colors[i], bold=True)
    add_text_box(slide, Inches(0.7), y + Inches(0.3), Inches(2.5), Inches(0.65),
                 content, font_size=9, color=DARK_TEXT)

# 右侧风险等级表
add_text_box(slide, Inches(3.8), Inches(1.5), Inches(4), Inches(0.35),
             "对公客户风险等级（四级）", font_size=14, color=NAVY, bold=True)
risk_table = [
    ("E1 成熟发展类", "业务经验丰富 · 制度健全 · 财务优良 · 信用风险极小", RGBColor(0x27, 0xAE, 0x60)),
    ("E2 稳健发展类", "具有业务经验 · 制度较健全 · 信用风险较小", MEDIUM_BLUE),
    ("E3 审慎关注类", "业务认知一般 · 制度相对完善 · 存在一定信用风险", ORANGE_ACCENT),
    ("E4 严格限制类", "存在重大风险隐患 · 禁止办理高风险业务", RED_ACCENT),
]
for i, (label, desc, color) in enumerate(risk_table):
    y = Inches(2.1) + Inches(0.7) * i
    add_rect(slide, Inches(3.8), y, Inches(0.35), Inches(0.55), color)
    add_text_box(slide, Inches(4.3), y, Inches(2.8), Inches(0.55),
                 label, font_size=12, color=color, bold=True)
    add_text_box(slide, Inches(7.0), y, Inches(5.5), Inches(0.55),
                 desc, font_size=10, color=DARK_TEXT)

# 底部：禁止准入 + 异常识别
add_gold_line(slide, Inches(3.8), Inches(5.0), Inches(8.8))

add_text_box(slide, Inches(3.8), Inches(5.15), Inches(8.8), Inches(0.3),
             "禁止准入客户清单", font_size=12, color=RED_ACCENT, bold=True)
ban_items = "身份造假 | 制裁命中 | 涉嫌洗钱/恐怖融资 | 司法冻结/失信被执行人 | 涉刑调查期间 | 被监管责令停业 | 破产清算 | 空壳公司"
add_text_box(slide, Inches(3.8), Inches(5.45), Inches(8.8), Inches(0.55),
             ban_items, font_size=9, color=DARK_TEXT)

add_text_box(slide, Inches(3.8), Inches(6.0), Inches(8.8), Inches(0.3),
             "异常客户识别规则（触发加强尽调 EDD）", font_size=12, color=ORANGE_ACCENT, bold=True)
edds = "客户信息与外部数据不一致 | 受益所有人涉及高风险国家/地区 | 注册地在FATF灰/黑名单地区 | 短期频繁变更法人/实控人 | 关联方涉及负面舆情 | 交易需求与经营规模不匹配"
add_text_box(slide, Inches(3.8), Inches(6.3), Inches(8.8), Inches(0.65),
             edds, font_size=9, color=DARK_TEXT)

# ==================== 第6页：端到端流程总览 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 2, "端到端业务操作全流程 —— 8 大环节总览")
add_page_number(slide, 6)

steps = [
    ("01", "客户申请\n预筛客", NAVY),
    ("02", "尽职调查\n(KYC)", MEDIUM_BLUE),
    ("03", "风险测评", GREEN_ACCENT),
    ("04", "适当性\n匹配", RGBColor(0x16, 0xA0, 0x85)),
    ("05", "准入审批", ORANGE_ACCENT),
    ("06", "签约开户", RGBColor(0x8E, 0x44, 0xAD)),
    ("07", "权限开通", RED_ACCENT),
    ("08", "交易/理财\n签约", RGBColor(0x2C, 0x3E, 0x50)),
]

for i, (num, label, color) in enumerate(steps):
    x = Inches(0.3) + Inches(1.55) * i
    y = Inches(1.6)
    # 圆形编号
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.4), y, Inches(0.7), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text_box(slide, x + Inches(0.4), y + Inches(0.1), Inches(0.7), Inches(0.5),
                 num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 步骤名
    add_text_box(slide, x, y + Inches(0.8), Inches(1.5), Inches(0.6),
                 label, font_size=10, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # 箭头（除最后一个）
    if i < len(steps) - 1:
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.25), y + Inches(0.2),
                                      Inches(0.35), Inches(0.25))
        arr.fill.solid()
        arr.fill.fore_color.rgb = GOLD
        arr.line.fill.background()

# 下半部分：关键节点说明
add_gold_line(slide, Inches(0.5), Inches(3.3), Inches(12.33))

key_points_left = [
    ("材料清单（对公）", [
        "① 营业执照(副本)  ② 公司章程",
        "③ 法定代表人身份证  ④ 授权委托书+被授权人身份证",
        "⑤ 近三年审计报告  ⑥ 企业征信报告",
        "⑦ KYC尽调表  ⑧ 风险测评问卷",
        "⑨ 代客外汇交易总协议  ⑩ 风险揭示书",
    ]),
]
for title, items in key_points_left:
    add_text_box(slide, Inches(0.5), Inches(3.5), Inches(6), Inches(0.3),
                 f"▎{title}", font_size=12, color=NAVY, bold=True)
    add_multi_text(slide, Inches(0.7), Inches(3.85), Inches(5.5), Inches(1.5),
                   [(item, 10, DARK_TEXT, False) for item in items])

# 审批层级表
add_text_box(slide, Inches(7), Inches(3.5), Inches(5.8), Inches(0.3),
             "▎审批层级矩阵", font_size=12, color=NAVY, bold=True)
approval_data = [
    ["客户类型", "业务范围", "审批层级"],
    ["个人客户", "I-III类", "支行行长"],
    ["个人客户", "IV-V类", "分行分管行长"],
    ["E1/E2对公", "I-III类", "支行→分行金融市场部"],
    ["E3对公", "I-II类", "支行→分行→分管行长"],
    ["超范围/豁免", "特殊审批", "总行金融市场部+合规部会签"],
]
for r, row in enumerate(approval_data):
    y = Inches(3.85) + Inches(0.35) * r
    for c, cell in enumerate(row):
        x = Inches(7.0) + Inches(1.9) * c
        is_header = (r == 0)
        add_text_box(slide, x, y, Inches(1.9), Inches(0.35),
                     cell, font_size=9, color=WHITE if is_header else DARK_TEXT,
                     bold=is_header, alignment=PP_ALIGN.CENTER)
        if is_header:
            add_rect(slide, x, y, Inches(1.85), Inches(0.35), NAVY)

# ==================== 第7页：KYC尽调与风险测评 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 2, "环节详解：尽调 → 测评 → 适当性匹配")
add_page_number(slide, 7)

# 左侧：KYC信息采集表
add_card(slide, Inches(0.35), Inches(1.5), Inches(4.0), Inches(2.5),
         "KYC 信息采集表（对公）", [
             ("基本信息", 10, NAVY, True),
             ("企业名称 | 成立时间 | 所属行业 | 注册地址 | 注册资本", 9, DARK_TEXT, False),
             ("", 6, DARK_TEXT, False),
             ("财务信息", 10, MEDIUM_BLUE, True),
             ("近三年营收 | 近三年净利润 | 资产负债率 | 主要结算银行", 9, DARK_TEXT, False),
             ("", 6, DARK_TEXT, False),
             ("尽调评估", 10, GREEN_ACCENT, True),
             ("衍生业务熟识度 (较高/一般/较低)", 9, DARK_TEXT, False),
             ("风险承受能力 (较高/一般/较低)", 9, DARK_TEXT, False),
             ("", 6, DARK_TEXT, False),
             ("会谈纪要", 10, ORANGE_ACCENT, True),
             ("会谈时间 | 参会人员 | 记录员 | 会谈主要内容", 9, DARK_TEXT, False),
             ("是否FT账户客户 | 实需交易背景描述", 9, DARK_TEXT, False),
         ], title_color=NAVY)

# 右上：风险测评评分维度
add_card(slide, Inches(4.7), Inches(1.5), Inches(4.1), Inches(2.5),
         "风险测评评分模型（对公）", [
             ("评估维度          权重", 10, NAVY, True),
             ("衍生业务熟识度    25%  较高80/一般50/较低20", 9, DARK_TEXT, False),
             ("风险承受能力       25%  较高80/一般50/较低20", 9, DARK_TEXT, False),
             ("财务状况             25%  营收+盈利能力+负债率", 9, DARK_TEXT, False),
             ("风控制度             15%  制度文件+人员配置评分", 9, DARK_TEXT, False),
             ("信用记录             10%  征信报告直接映射评分", 9, DARK_TEXT, False),
             ("", 6, DARK_TEXT, False),
             ("总分映射：≥85→E1  65-84→E2  40-64→E3  <40→E4", 10, ORANGE_ACCENT, True),
         ], title_color=MEDIUM_BLUE)

# 右下：适当性匹配表
add_card(slide, Inches(9.15), Inches(1.5), Inches(3.85), Inches(2.5),
         "适当性匹配规则", [
             ("产品分类", 10, NAVY, True),
             ("I类(低风险)   即期结售汇/外汇买卖", 9, DARK_TEXT, False),
             ("II类(中低)     远期结售汇/外汇买卖/买入期权", 9, DARK_TEXT, False),
             ("III类(中)       复杂期权/利率互换", 9, DARK_TEXT, False),
             ("IV类(中高)     结构化远期/奇异期权/CCS", 9, DARK_TEXT, False),
             ("V类(高)         杠杆产品/非标复杂组合", 9, DARK_TEXT, False),
             ("", 6, DARK_TEXT, False),
             ("核心规则：客户等级 ≥ 产品等级", 10, RED_ACCENT, True),
             ("超风险销售 → 系统拦截 → 须特殊审批", 10, RED_ACCENT, False),
         ], title_color=GREEN_ACCENT)

# 底部：适当性匹配矩阵
add_text_box(slide, Inches(0.35), Inches(4.3), Inches(6), Inches(0.3),
             "▎客户-产品适当性匹配矩阵", font_size=12, color=NAVY, bold=True)

matrix_headers = ["客户等级", "I类 低风险", "II类 中低风险", "III类 中风险", "IV类 中高风险", "V类 高风险"]
matrix_data = [
    ["E1 成熟发展类", "✓ 准入", "✓ 准入", "✓ 准入", "✓ 准入", "△ 特殊审批"],
    ["E2 稳健发展类", "✓ 准入", "✓ 准入", "✓ 准入", "△ 特殊审批", "✕ 禁止"],
    ["E3 审慎关注类", "✓ 准入", "✓ 准入", "△ 特殊审批", "✕ 禁止", "✕ 禁止"],
    ["E4 严格限制类", "△ 特殊审批", "✕ 禁止", "✕ 禁止", "✕ 禁止", "✕ 禁止"],
]

for r, row in enumerate([matrix_headers] + matrix_data):
    y = Inches(4.7) + Inches(0.38) * r
    for c, cell in enumerate(row):
        is_header = (r == 0)
        widths = [1.5, 1.6, 1.6, 1.6, 1.6, 1.6]
        x = Inches(0.35) + sum(widths[:c])
        if is_header:
            add_rect(slide, x, y, Inches(widths[c]), Inches(0.38), NAVY)
            add_text_box(slide, x, y, Inches(widths[c]), Inches(0.38),
                         cell, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        else:
            bg = LIGHT_GRAY if c == 0 else (RGBColor(0xE8, 0xF8, 0xF0) if "✓" in cell else
                  (RGBColor(0xFD, 0xED, 0xEC) if "✕" in cell else RGBColor(0xFE, 0xF5, 0xE7)))
            add_rect(slide, x, y, Inches(widths[c]), Inches(0.38), bg)
            txt_color = GREEN_ACCENT if "✓" in cell else (RED_ACCENT if "✕" in cell else (ORANGE_ACCENT if "△" in cell else DARK_TEXT))
            add_text_box(slide, x, y, Inches(widths[c]), Inches(0.38),
                         cell, font_size=9, color=txt_color, bold=True, alignment=PP_ALIGN.CENTER)

# ==================== 第8页：签约开户+权限开通 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 2, "环节详解：审批 → 签约 → 权限开通 → 交易")
add_page_number(slide, 8)

# 线下签约流程
add_card(slide, Inches(0.35), Inches(1.5), Inches(4.2), Inches(3.0),
         "线下签约流程", [
             ("审批通过 → 客户经理通知客户", 11, NAVY, True),
             ("", 4, NAVY, False),
             ("签署文件清单：", 10, MEDIUM_BLUE, True),
             ("① 《代客外汇交易总协议》（主协议）", 9, DARK_TEXT, False),
             ("② 《风险揭示书》", 9, DARK_TEXT, False),
             ("③ 《交易确认书模板》", 9, DARK_TEXT, False),
             ("④ 《客户适当性匹配确认书》", 9, DARK_TEXT, False),
             ("⑤ 《企业授权委托书》（对公）", 9, DARK_TEXT, False),
             ("⑥ 《业务合规承诺函》", 9, DARK_TEXT, False),
             ("", 4, DARK_TEXT, False),
             ("加盖公章（对公）或签名确认（个人）", 9, DARK_TEXT, False),
             ("客户经理将签署文件扫描上传至系统", 9, DARK_TEXT, False),
         ], title_color=NAVY)

# 线上签约流程
add_card(slide, Inches(4.9), Inches(1.5), Inches(3.9), Inches(3.0),
         "线上签约流程", [
             ("客户登录企业/个人网银", 11, NAVY, True),
             ("", 4, NAVY, False),
             ("操作步骤：", 10, MEDIUM_BLUE, True),
             ("① 进入「协议签约」模块", 9, DARK_TEXT, False),
             ("② 阅读协议全文（强制阅读≥30秒）", 9, DARK_TEXT, False),
             ("③ U-Key数字签名 / 人脸识别确认", 9, DARK_TEXT, False),
             ("④ 短信验证码复核", 9, DARK_TEXT, False),
             ("⑤ 系统生成签约记录，自动归档", 9, DARK_TEXT, False),
             ("", 4, DARK_TEXT, False),
             ("合规要求：", 10, RED_ACCENT, True),
             ("线上签约需双因素认证", 9, DARK_TEXT, False),
             ("全程留痕（浏览时长+电子签名时间戳）", 9, DARK_TEXT, False),
             ("个人理财：全程双录（录音+录像）", 9, DARK_TEXT, False),
         ], title_color=MEDIUM_BLUE)

# 权限开通
add_card(slide, Inches(9.15), Inches(1.5), Inches(3.85), Inches(3.0),
         "权限开通 & 额度管理", [
             ("签约完成 → 系统自动赋权", 11, NAVY, True),
             ("", 4, NAVY, False),
             ("权限维度（四维矩阵）：", 10, MEDIUM_BLUE, True),
             ("• 交易品种：按客户等级匹配", 9, DARK_TEXT, False),
             ("• 交易额度：单笔限额+累计限额+年交割量", 9, DARK_TEXT, False),
             ("• 交易渠道：网银/柜面/API直连", 9, DARK_TEXT, False),
             ("• 保证金/授信：按信用等级差异化", 9, DARK_TEXT, False),
             ("", 4, DARK_TEXT, False),
             ("年建议交割量公式：", 10, GREEN_ACCENT, True),
             ("大中型 = 近12月实需×(1+增长率)-已签约", 9, DARK_TEXT, False),
             ("小微 = 近12月实需×(1+增长率)×80%-已签约", 9, DARK_TEXT, False),
         ], title_color=GREEN_ACCENT)

# 底部异常场景
add_gold_line(slide, Inches(0.35), Inches(4.8), Inches(12.63))

add_text_box(slide, Inches(0.35), Inches(5.0), Inches(12), Inches(0.3),
             "▎业务异常场景处理", font_size=14, color=NAVY, bold=True)

scenarios = [
    ("拒绝准入", "反洗钱不通过/材料不符/审批不通过", "通知客户→记录原因→6个月冷却期→可重新申请", RED_ACCENT),
    ("暂缓准入", "材料不完整/信息待核实/风险等级待确认", "标记「待补充」→暂缓≤30天→超期退回→重新申请", ORANGE_ACCENT),
    ("白名单准入", "总行级战略客户/大型国企/优质上市公司", "总行发起→白名单维护→简化尽调→快速通道", GREEN_ACCENT),
    ("特殊审批", "客户等级<产品等级/超额度/超范围品种", "发起特批→分行→总行→合规会签→设有效期", MEDIUM_BLUE),
]

for i, (title, trigger, action, color) in enumerate(scenarios):
    x = Inches(0.35) + Inches(3.2) * i
    add_rect(slide, x, Inches(5.4), Inches(2.95), Pt(4), color)
    add_text_box(slide, x, Inches(5.5), Inches(2.95), Inches(0.35),
                 title, font_size=12, color=color, bold=True)
    add_text_box(slide, x, Inches(5.85), Inches(2.95), Inches(0.35),
                 f"触发：{trigger}", font_size=8, color=GRAY)
    add_text_box(slide, x, Inches(6.15), Inches(2.95), Inches(0.55),
                 f"处置：{action}", font_size=8, color=DARK_TEXT)

# ==================== 第9页：系统架构 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 3, "系统实现方案 —— 整体架构设计")
add_page_number(slide, 9)

# 架构分层
layers = [
    ("客户接入层", NAVY, Inches(1.6), [
        "个人网银", "企业网银", "柜面终端", "API直连",
    ]),
    ("客户准入管理平台（准入中台）", GOLD, Inches(2.5), [
        "客户信息采集", "尽调录入管理", "风险测评", "准入审批流",
        "黑名单校验", "适当性匹配引擎", "权限管理", "日志审计",
    ]),
    ("系统集成层", MEDIUM_BLUE, Inches(4.05), [
        "核心业务系统\n(CBS)", "CRM\n系统", "反洗钱\n(AML)", "风控系统\n(RISK)",
        "ECIF\n客户信息中心", "额度管理\n系统", "二代征信\n系统", "工商信息\n联网核查",
        "身份认证\n平台", "电子签章\n平台", "电子档案\n平台", "监管报送\n平台",
    ]),
]

for title, color, y, modules in layers:
    # 层标签
    label = add_rect(slide, Inches(0.35), y, Inches(2.0), Inches(0.65), color)
    add_text_box(slide, Inches(0.35), y + Inches(0.1), Inches(2.0), Inches(0.45),
                 title, font_size=11, color=WHITE if color != GOLD else NAVY_DARK,
                 bold=True, alignment=PP_ALIGN.CENTER)
    # 模块
    mw = Inches(2.4) if len(modules) <= 4 else Inches(1.8)
    for i, mod in enumerate(modules):
        mx = Inches(2.6) + (mw + Inches(0.1)) * i
        mod_bg = RGBColor(0xF0, 0xF4, 0xF8) if color != GOLD else RGBColor(0x3A, 0x4A, 0x6A)
        mod_color = DARK_TEXT if color != GOLD else WHITE
        add_rect(slide, mx, y + Inches(0.05), mw, Inches(0.55), mod_bg)
        add_text_box(slide, mx, y + Inches(0.08), mw, Inches(0.5),
                     mod, font_size=9, color=mod_color, bold=True, alignment=PP_ALIGN.CENTER)

# 集成箭头
for i in range(2):
    y = Inches(2.25) + Inches(1.55) * i
    arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(1.1), y, Inches(0.3), Inches(0.25))
    arr.fill.solid()
    arr.fill.fore_color.rgb = GOLD
    arr.line.fill.background()

# 底部：系统交互流向
add_gold_line(slide, Inches(0.35), Inches(5.0), Inches(12.63))

add_text_box(slide, Inches(0.35), Inches(5.15), Inches(12), Inches(0.3),
             "▎关键数据同步", font_size=14, color=NAVY, bold=True)

sync_data = [
    ["数据", "源系统", "目标系统", "同步方式", "频率"],
    ["客户基础信息", "ECIF", "准入平台", "ESB服务总线", "实时"],
    ["工商信息", "国家企业信用信息公示系统", "准入平台", "API调用", "查询时实时"],
    ["黑名单状态", "AML系统", "准入平台", "API调用", "各准入节点实时"],
    ["征信信息", "二代征信系统", "准入平台", "批量文件", "T+1"],
    ["准入结果", "准入平台", "核心系统(CBS)", "ESB服务总线", "审批通过后实时"],
    ["交易权限", "准入平台", "交易系统", "ESB服务总线", "权限变更后实时"],
]

for r, row in enumerate(sync_data):
    y = Inches(5.5) + Inches(0.28) * r
    is_header = (r == 0)
    widths = [2.8, 2.8, 2.8, 2.2, 2.0]
    for c, cell in enumerate(row):
        x = Inches(0.35) + sum(widths[:c])
        if is_header:
            add_rect(slide, x, y, Inches(widths[c]), Inches(0.28), NAVY)
            add_text_box(slide, x, y, Inches(widths[c]), Inches(0.28),
                         cell, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        else:
            add_text_box(slide, x, y, Inches(widths[c]), Inches(0.28),
                         cell, font_size=8, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# ==================== 第10页：核心功能模块 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 3, "系统功能模块详解 —— 8 大核心模块")
add_page_number(slide, 10)

modules_list = [
    ("客户信息采集", NAVY, "ECIF对接 · OCR识别\n身份证/营业执照联网核查\n受益所有人自动穿透\n九要素完整性校验"),
    ("尽调录入管理", MEDIUM_BLUE, "结构化KYC表单模板\n影像资料上传+OCR识别\nFT账户特殊逻辑处理\n行业分类自动匹配"),
    ("风险测评", GREEN_ACCENT, "线上/线下双模式\n自动评分→等级映射\n历史测评追踪对比\n到期30/7/1天提醒"),
    ("准入校验引擎", ORANGE_ACCENT, "Drools规则引擎驱动\n8大校验规则实时执行\n制裁名单API实时调用\n规则在线热更新"),
    ("审批工作流", RGBColor(0x8E, 0x44, 0xAD), "Flowable工作流引擎\n5级审批可配置\n48h超时自动升签\n审批全链路回溯"),
    ("权限管理", RED_ACCENT, "4维权限矩阵模型\n审批后自动赋权\n权限变更审批流程\n实时生效无需刷新"),
    ("黑名单校验", RGBColor(0x2C, 0x3E, 0x50), "对接AML反洗钱系统\n多源名单实时筛查\n准入5节点触发校验\n结果缓存24h"),
    ("日志审计", RGBColor(0x16, 0xA0, 0x85), "全链路不可篡改日志\n操作前后数据变更详情\n≥5年日志保留\n准入统计+异常报表"),
]

for i, (title, color, desc) in enumerate(modules_list):
    col = i % 4
    row = i // 4
    x = Inches(0.35) + Inches(3.2) * col
    y = Inches(1.5) + Inches(2.8) * row
    # 卡片
    add_rect(slide, x, y, Inches(2.95), Inches(1.55), WHITE)
    add_rect(slide, x, y, Inches(2.95), Pt(4), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.65), Inches(0.35),
                 f"{i+1}. {title}", font_size=13, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.55), Inches(2.65), Inches(1.0),
                 desc, font_size=9, color=DARK_TEXT)

# ==================== 第11页：规则引擎 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 3, "关键规则引擎 —— 准入校验 & 适当性匹配")
add_page_number(slide, 11)

# 准入条件自动校验引擎
add_text_box(slide, Inches(0.35), Inches(1.5), Inches(6), Inches(0.3),
             "▎准入条件自动校验引擎（Drools 规则引擎）", font_size=14, color=NAVY, bold=True)

checks = [
    ("黑名单校验", "客户姓名/证件号/法人/实控人 in AML黑名单", "直接拒绝准入", RED_ACCENT),
    ("制裁名单", "命中联合国/OFAC/央行制裁名单", "直接拒绝+上报合规", RED_ACCENT),
    ("PEPS筛查", "客户为政治公众人物或密切关系人", "触发EDD流程", ORANGE_ACCENT),
    ("行业禁入", "客户所属行业 in 禁入清单", "拒绝准入", RED_ACCENT),
    ("证件有效期", "证件已过期或距到期不足30天", "提示更新后申请", ORANGE_ACCENT),
    ("KYC完整性", "必填字段未完整填写", "退回补充", GRAY),
    ("受益所有人", "受益所有人未穿透至自然人", "退回补充", ORANGE_ACCENT),
    ("适当性匹配", "客户风险等级 < 申请产品风险等级", "拒绝开通品种权限", RED_ACCENT),
]

for i, (name, rule, action, color) in enumerate(checks):
    y = Inches(1.95) + Inches(0.5) * i
    add_rect(slide, Inches(0.35), y, Inches(0.25), Inches(0.38), color)
    add_text_box(slide, Inches(0.75), y, Inches(2.2), Inches(0.38),
                 name, font_size=10, color=color, bold=True)
    add_text_box(slide, Inches(2.9), y, Inches(4.5), Inches(0.38),
                 f"WHEN: {rule}", font_size=9, color=DARK_TEXT)
    add_text_box(slide, Inches(7.4), y, Inches(3.0), Inches(0.38),
                 f"THEN: {action}", font_size=9, color=color, bold=True)

# 右侧：适当性匹配引擎
add_text_box(slide, Inches(0.35), Inches(6.1), Inches(12), Inches(0.3),
             "▎适当性匹配自动判断引擎", font_size=14, color=NAVY, bold=True)

add_text_box(slide, Inches(0.35), Inches(6.45), Inches(5.5), Inches(0.25),
             "输入 → 客户风险等级 + 产品风险等级 + 产品类别", font_size=10, color=DARK_TEXT)
add_text_box(slide, Inches(0.35), Inches(6.7), Inches(5.5), Inches(0.25),
             "输出 → 匹配结果(通过/不通过) + 可交易品种清单", font_size=10, color=DARK_TEXT)

# 右侧特别规则
add_text_box(slide, Inches(7.5), Inches(6.1), Inches(5.5), Inches(0.3),
             "▎特别规则（不可绕过）", font_size=12, color=RED_ACCENT, bold=True)
add_multi_text(slide, Inches(7.5), Inches(6.45), Inches(5.5), Inches(0.8), [
    ("• 杠杆类产品(V类)：仅C5/E1客户可准入", 10, DARK_TEXT, False),
    ("• 卖出期权：仅风险等级 ≥ C4/E2 可准入", 10, DARK_TEXT, False),
    ("• FT账户客户：额外校验跨境交易合规要求", 10, DARK_TEXT, False),
])

# ==================== 第12页：系统数据流转 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 3, "数据流转 & 接口设计")
add_page_number(slide, 12)

# 数据流转图
add_text_box(slide, Inches(0.35), Inches(1.5), Inches(6), Inches(0.3),
             "▎客户数据流转架构", font_size=14, color=NAVY, bold=True)

# 源系统
sources = ["工商信息\n联网核查", "ECIF\n客户信息中心", "CRM\n系统"]
for i, src in enumerate(sources):
    x = Inches(0.35) + Inches(1.8) * i
    add_rect(slide, x, Inches(2.1), Inches(1.6), Inches(0.65), MEDIUM_BLUE)
    add_text_box(slide, x, Inches(2.15), Inches(1.6), Inches(0.55),
                 src, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 箭头汇聚
for i in range(3):
    x = Inches(0.35) + Inches(1.8) * i + Inches(0.65)
    arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, Inches(2.8), Inches(0.3), Inches(0.3))
    arr.fill.solid()
    arr.fill.fore_color.rgb = GOLD
    arr.line.fill.background()

# 准入中台
add_rect(slide, Inches(0.6), Inches(3.2), Inches(4.5), Inches(0.6), NAVY)
add_text_box(slide, Inches(0.6), Inches(3.3), Inches(4.5), Inches(0.4),
             "客户准入管理平台（准入中台）", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 分散箭头
for i in range(4):
    x = Inches(0.6) + Inches(1.2) * i + Inches(0.35)
    arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, Inches(3.85), Inches(0.3), Inches(0.3))
    arr.fill.solid()
    arr.fill.fore_color.rgb = GOLD
    arr.line.fill.background()

# 目标系统
targets = ["反洗钱系统\n(AML)", "风控系统\n(RISK)", "核心系统\n(CBS)", "电子档案\n平台"]
for i, tgt in enumerate(targets):
    x = Inches(0.35) + Inches(1.3) * i
    add_rect(slide, x, Inches(4.25), Inches(1.15), Inches(0.65), LIGHT_GRAY)
    add_text_box(slide, x, Inches(4.28), Inches(1.15), Inches(0.55),
                 tgt, font_size=9, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

# 右侧-接口设计
add_text_box(slide, Inches(6.0), Inches(1.5), Inches(7), Inches(0.3),
             "▎接口设计要点", font_size=14, color=NAVY, bold=True)

api_items = [
    ("协议标准", "内部系统：RESTful API（JSON）\n存量系统过渡：Web Service（XML）"),
    ("安全保障", "API网关统一鉴权\n传输层TLS 1.2+加密\n报文敏感字段脱敏（身份证/手机号）"),
    ("幂等设计", "所有写操作接口支持幂等\n防止重复提交（redis分布式锁）\n超时重试机制（3次+指数退避）"),
    ("监控告警", "接口调用量/延迟/错误率实时监控\nESB同步状态监控面板\n同步失败自动重试+钉钉/短信告警"),
]

for i, (title, content) in enumerate(api_items):
    y = Inches(1.95) + Inches(1.2) * i
    add_rect(slide, Inches(6.0), y, Inches(0.25), Inches(0.9), GOLD)
    add_text_box(slide, Inches(6.4), y, Inches(2.0), Inches(0.3),
                 title, font_size=11, color=NAVY, bold=True)
    add_text_box(slide, Inches(6.4), y + Inches(0.3), Inches(5.5), Inches(0.65),
                 content, font_size=9, color=DARK_TEXT)

# ==================== 第13页：风险全景 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 4, "风险点与管控措施 —— 三大风险全景")
add_page_number(slide, 13)

risks = [
    ("合规风险", RED_ACCENT, [
        ("监管处罚风险", "向不满足准入条件客户\n开通业务，被监管责令\n整改/罚款", "高"),
        ("适当性违规", "向低风险承受能力客户\n销售高风险产品，违反\n适当性管理要求", "高"),
        ("反洗钱违规", "KYC尽调不到位、未准确\n识别受益所有人、未筛查\n制裁名单", "高"),
        ("越权审批", "审批人员超越权限审批\n导致不合规客户准入", "中"),
        ("材料造假", "客户经理协助客户\n伪造准入材料", "中"),
    ]),
    ("操作风险", ORANGE_ACCENT, [
        ("信息录入错误", "客户经理录入信息错误\n导致风险评估偏差", "中"),
        ("审批超时", "审批人未及时处理\n导致流程阻塞", "中"),
        ("权限配置错误", "系统或人工错误配置\n超额/超范围授权", "中"),
        ("数据同步失败", "准入结果未成功同步\n客户无法交易", "低"),
        ("KYC过期未处理", "到期后未自动关闭\n交易权限", "中"),
    ]),
    ("洗钱风险", RGBColor(0x8E, 0x44, 0xAD), [
        ("壳公司准入", "空壳公司伪造贸易背景\n用于非法资金跨境转移", "高"),
        ("受益所有人隐匿", "多层股权结构隐匿实控人\n规避制裁筛查", "高"),
        ("虚假贸易背景", "伪造合同/发票构造虚假\n实需，实现资本外逃", "高"),
        ("化整为零", "多个关联主体分散交易\n规避大额交易监控", "中"),
        ("跨境洗钱", "利用复杂外汇交易\n掩盖资金来源", "高"),
    ]),
]

for col_idx, (cat_name, cat_color, items) in enumerate(risks):
    x = Inches(0.35) + Inches(4.3) * col_idx
    # 分类标题
    add_rect(slide, x, Inches(1.5), Inches(4.0), Inches(0.45), cat_color)
    add_text_box(slide, x, Inches(1.5), Inches(4.0), Inches(0.45),
                 cat_name, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    for i, (name, desc, level) in enumerate(items):
        y = Inches(2.1) + Inches(1.05) * i
        lvl_color = RED_ACCENT if level == "高" else (ORANGE_ACCENT if level == "中" else GREEN_ACCENT)
        add_rect(slide, x, y, Inches(0.25), Inches(0.85), lvl_color)
        add_text_box(slide, x + Inches(0.35), y, Inches(1.3), Inches(0.85),
                     name, font_size=10, color=DARK_TEXT, bold=True)
        add_text_box(slide, x + Inches(1.7), y, Inches(1.6), Inches(0.85),
                     desc, font_size=8, color=GRAY)
        add_text_box(slide, x + Inches(3.35), y, Inches(0.5), Inches(0.85),
                     f"[{level}]", font_size=10, color=lvl_color, bold=True)

# ==================== 第14页：管控措施 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 4, "风险管控体系 —— 事前·事中·事后三层防控")
add_page_number(slide, 14)

# 三层
controls = [
    ("事前预防", NAVY, [
        "身份证/营业执照联网核查\n人脸识别实名比对",
        "实时对接AML系统\n5节点触发黑名单筛查",
        "PEPS政治人物自动标记\n触发加强尽调EDD",
        "行业禁入清单自动拦截\n适当性匹配自动校验",
    ]),
    ("事中监控", MEDIUM_BLUE, [
        "交易权限硬控制\n超权限交易系统绝对拦截",
        "额度累控实时计算\n超建议交割量预警",
        "异常交易检测\n时间/频率/对手异常",
        "双录AI质检(个人理财)\n自动检测误导销售用语",
    ]),
    ("事后审计", GREEN_ACCENT, [
        "全链路不可篡改日志\n5分钟还原准入全流程",
        "按日/周/月生成\n准入合规报表",
        "风险事件回溯追踪\n关联交易链还原",
        "定期合规抽查\n存量客户年度复评",
    ]),
]

for i, (title, color, items) in enumerate(controls):
    x = Inches(0.35) + Inches(4.3) * i
    add_rect(slide, x, Inches(1.5), Inches(4.0), Inches(0.6), color)
    add_text_box(slide, x, Inches(1.55), Inches(4.0), Inches(0.5),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.2), Inches(2.3) + Inches(1.1) * j,
                     Inches(3.6), Inches(1.0),
                     f"▸  {item}", font_size=10, color=DARK_TEXT)

# 底部监控表
add_gold_line(slide, Inches(0.35), Inches(5.65), Inches(12.63))
add_text_box(slide, Inches(0.35), Inches(5.75), Inches(12), Inches(0.3),
             "▎监控预警矩阵", font_size=14, color=NAVY, bold=True)

alerts = [
    ["类别", "监控项", "监控方式", "告警条件"],
    ["系统级", "AML系统可用性", "心跳检测", "连续3次无响应→告警"],
    ["系统级", "审批流程阻塞", "待办任务超时扫描", "超过48h未处理→告警"],
    ["系统级", "数据同步失败", "ESB同步状态监控", "同步失败→即时告警"],
    ["系统级", "KYC到期", "定时扫描任务", "到期前30/7/1天→提醒"],
    ["业务级", "高风险客户集中准入", "准入数据实时分析", "单日E3/E4超阈值→通知合规"],
    ["业务级", "异常审批模式", "审批行为分析模型", "同审批人连续快速通过→标记"],
    ["业务级", "客户信息频繁变更", "变更频率监控", "30天内法人/实控人变更≥2次→EDD"],
    ["业务级", "交易量突增", "交易量环比分析", "月交易量>建议交割量150%→核实"],
]

for r, row in enumerate(alerts):
    y = Inches(6.1) + Inches(0.16) * r
    is_header = (r == 0)
    widths = [1.2, 3.0, 3.0, 5.5]
    for c, cell in enumerate(row):
        x = Inches(0.35) + sum(widths[:c])
        if is_header:
            add_rect(slide, x, y, Inches(widths[c]), Inches(0.16), NAVY)
            add_text_box(slide, x, y, Inches(widths[c]), Inches(0.16),
                         cell, font_size=7, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        else:
            add_text_box(slide, x, y, Inches(widths[c]), Inches(0.16),
                         cell, font_size=7, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# ==================== 第15页：总结 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY_DARK)

add_rect(slide, Inches(0.8), Inches(1.8), Pt(6), Inches(4.0), GOLD)

add_text_box(slide, Inches(1.3), Inches(2.0), Inches(10), Inches(0.8),
             "方案总结", font_size=30, color=WHITE, bold=True)
add_gold_line(slide, Inches(1.3), Inches(2.8), Inches(3))

summary_items = [
    ("合规先行", "严格遵循8项核心监管法规，建立五维客户准入模型，覆盖个人+对公全客群"),
    ("流程闭环", "8大环节端到端贯通，6级审批矩阵，4类异常场景全覆盖，材料清单+合规校验点逐环落地"),
    ("系统落地", "准入中台架构 + 8大功能模块 + 3大规则引擎 + 12系统对接，提供完整技术方案"),
    ("风控兜底", "三大风险15个风险点逐一识别，事前·事中·事后三层管控，系统+业务双维度预警"),
]

for i, (title, desc) in enumerate(summary_items):
    y = Inches(3.2) + Inches(1.0) * i
    add_text_box(slide, Inches(1.3), y, Inches(1.8), Inches(0.4),
                 title, font_size=18, color=GOLD, bold=True)
    add_text_box(slide, Inches(3.2), y + Inches(0.05), Inches(8), Inches(0.4),
                 desc, font_size=12, color=LIGHT_GRAY)

add_text_box(slide, Inches(1.3), Inches(7.0), Inches(5), Inches(0.3),
             "金融市场部 · 合规部 · 风险管理部 · 信息科技部  |  2025年5月  |  V1.0  |  内部限制",
             font_size=9, color=GRAY)
add_page_number(slide, 15)

# ==================== 第16页：尾页 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY_DARK)

add_text_box(slide, Inches(0), Inches(2.5), Inches(13.33), Inches(1.2),
             "THANK YOU", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_gold_line(slide, Inches(5.3), Inches(3.7), Inches(2.67))
add_text_box(slide, Inches(0), Inches(4.0), Inches(13.33), Inches(0.8),
             "商业银行代客外汇交易、理财代销业务 · 客户准入全流程及系统实现方案",
             font_size=14, color=GOLD_LIGHT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(4.8), Inches(13.33), Inches(0.6),
             "内部资料 · 请勿外传", font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 16)

# ==================== 保存 ====================
output_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                           'docs', '客户准入全流程及系统实现方案.pptx')
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
